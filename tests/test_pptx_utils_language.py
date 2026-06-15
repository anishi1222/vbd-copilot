"""Regression tests for the language / East Asian font support added to
``skills/pptx-generator/pptx_utils.py``.

These tests cover the ``set_language()`` API and verify that runs created
through public helpers carry the expected ``<a:latin>`` and ``<a:ea>``
typeface elements when Japanese output is selected.
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest
from pptx.oxml.ns import qn
from pptx.util import Inches

# Load pptx_utils from the (hyphenated) skill directory so it can be imported
# as a module without modifying sys.path globally for the rest of the suite.
_SPEC = importlib.util.spec_from_file_location(
    "_pptx_utils_under_test",
    Path(__file__).parent.parent / "skills" / "pptx-generator" / "pptx_utils.py",
)
pu = importlib.util.module_from_spec(_SPEC)
sys.modules["_pptx_utils_under_test"] = pu
_SPEC.loader.exec_module(pu)


@pytest.fixture(autouse=True)
def _reset_language():
    """Restore the default language after every test so module state never leaks."""
    yield
    pu.set_language("en")


def _rPr(run):
    return run._r.find(qn("a:rPr"))


def _typeface(rPr, child_tag: str) -> str | None:
    el = rPr.find(qn(child_tag))
    return el.get("typeface") if el is not None else None


# ── set_language API ─────────────────────────────────────────────────────────


class TestSetLanguage:
    def test_default_language_is_en(self):
        assert pu.active_language() == "en"

    def test_set_language_ja(self):
        pu.set_language("ja")
        assert pu.active_language() == "ja"

    def test_set_language_back_to_en(self):
        pu.set_language("ja")
        pu.set_language("en")
        assert pu.active_language() == "en"

    def test_invalid_language_raises(self):
        with pytest.raises(ValueError, match="Unknown language"):
            pu.set_language("zh")

    def test_invalid_language_does_not_change_state(self):
        pu.set_language("ja")
        with pytest.raises(ValueError):
            pu.set_language("fr")
        assert pu.active_language() == "ja"


# ── _ja_ea_typeface_for mapping ──────────────────────────────────────────────


class TestEaTypefaceMapping:
    def test_segoe_ui_maps_to_yu_gothic_ui(self):
        assert pu._ja_ea_typeface_for(pu.FONT_FAMILY) == "Yu Gothic UI"

    def test_segoe_ui_semibold_maps_to_yu_gothic_ui(self):
        # Bold weight is carried by run.font.bold; the EA family stays regular.
        assert pu._ja_ea_typeface_for(pu.FONT_SEMIBOLD) == "Yu Gothic UI"

    def test_segoe_ui_light_maps_to_yu_gothic_ui_light(self):
        assert pu._ja_ea_typeface_for(pu.FONT_LIGHT) == "Yu Gothic UI Light"

    def test_unknown_typeface_falls_back_to_yu_gothic_ui(self):
        assert pu._ja_ea_typeface_for("Some Other Font") == "Yu Gothic UI"


# ── _apply_run_font behaviour (unit, no real slide) ──────────────────────────


def _make_lone_run():
    """Create an isolated python-pptx Run wrapper for typeface assertions."""
    prs = pu.create_presentation()
    slide = pu.new_blank_slide(prs)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "X"
    return run


class TestApplyRunFont:
    def test_en_sets_only_latin(self):
        run = _make_lone_run()
        pu._apply_run_font(run, pu.FONT_FAMILY)
        rPr = _rPr(run)
        assert _typeface(rPr, "a:latin") == "Segoe UI"
        assert _typeface(rPr, "a:ea") is None

    def test_ja_sets_latin_and_ea(self):
        pu.set_language("ja")
        run = _make_lone_run()
        pu._apply_run_font(run, pu.FONT_FAMILY)
        rPr = _rPr(run)
        assert _typeface(rPr, "a:latin") == "Segoe UI"
        assert _typeface(rPr, "a:ea") == "Yu Gothic UI"

    def test_ja_with_set_ea_false_skips_ea(self):
        pu.set_language("ja")
        run = _make_lone_run()
        pu._apply_run_font(run, pu.FONT_MONO, set_ea=False)
        rPr = _rPr(run)
        assert _typeface(rPr, "a:latin") == "Courier New"
        assert _typeface(rPr, "a:ea") is None

    def test_ja_re_applying_does_not_duplicate_ea(self):
        pu.set_language("ja")
        run = _make_lone_run()
        pu._apply_run_font(run, pu.FONT_FAMILY)
        pu._apply_run_font(run, pu.FONT_SEMIBOLD)
        rPr = _rPr(run)
        eas = rPr.findall(qn("a:ea"))
        assert len(eas) == 1
        assert eas[0].get("typeface") == "Yu Gothic UI"

    def test_ja_light_typeface_uses_light_ea(self):
        pu.set_language("ja")
        run = _make_lone_run()
        pu._apply_run_font(run, pu.FONT_LIGHT)
        rPr = _rPr(run)
        assert _typeface(rPr, "a:ea") == "Yu Gothic UI Light"


# ── Public helpers honour the language setting ───────────────────────────────


class TestPublicHelpersInJapaneseMode:
    def test_add_textbox_writes_ea_in_ja_mode(self):
        pu.set_language("ja")
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        tb = pu.add_textbox(
            slide,
            "\u3053\u3093\u306b\u3061\u306f Azure",
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(1),
        )
        run = tb.text_frame.paragraphs[0].runs[0]
        rPr = _rPr(run)
        assert _typeface(rPr, "a:ea") == "Yu Gothic UI"

    def test_add_textbox_no_ea_in_en_mode(self):
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        tb = pu.add_textbox(
            slide,
            "Hello world",
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(1),
        )
        run = tb.text_frame.paragraphs[0].runs[0]
        rPr = _rPr(run)
        assert _typeface(rPr, "a:ea") is None

    def test_bold_markup_keeps_ea_for_each_run(self):
        pu.set_language("ja")
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        tb = pu.add_textbox(
            slide,
            "Use **Copilot** here",
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(1),
        )
        runs = tb.text_frame.paragraphs[0].runs
        assert len(runs) == 3
        for r in runs:
            rPr = _rPr(r)
            assert _typeface(rPr, "a:ea") == "Yu Gothic UI", (
                f"Run {r.text!r} missing East Asian typeface"
            )
        # The bold middle run uses Segoe UI Semibold for Latin.
        assert _typeface(_rPr(runs[1]), "a:latin") == "Segoe UI Semibold"
        assert runs[1].font.bold is True

    def test_code_block_runs_have_no_ea(self):
        """Monospaced code blocks must keep alignment intact regardless of language."""
        pu.set_language("ja")
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        pu.add_code_block(
            slide,
            "def foo():\n    return 1",
            Inches(1),
            Inches(1),
            Inches(8),
        )

        def _walk_runs(shape):
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for child in shape.shapes:
                    yield from _walk_runs(child)
            elif shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    yield from p.runs

        code_runs = [
            r
            for r in _walk_runs(slide.shapes[0])
            if _typeface(_rPr(r), "a:latin") == "Courier New"
        ]
        assert code_runs, "no Courier New runs found in code block"
        for r in code_runs:
            assert _typeface(_rPr(r), "a:ea") is None

    def test_gradient_textbox_preserves_ooxml_child_order(self):
        """`<a:gradFill>` must precede `<a:latin>` and `<a:ea>` per OOXML schema."""
        pu.set_language("ja")
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        tb = pu.add_gradient_textbox(
            slide,
            "\u30b0\u30e9\u30c7\u30fc\u30b7\u30e7\u30f3 Azure",
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(1),
            pu.MS_BLUE,
            pu.MS_LIGHT_BLUE,
            font_size=28,
            bold=True,
        )
        run = tb.text_frame.paragraphs[0].runs[0]
        rPr = _rPr(run)
        tags = [child.tag.split("}", 1)[-1] for child in rPr]
        # gradFill must come before latin which must come before ea
        assert tags.index("gradFill") < tags.index("latin") < tags.index("ea")
        assert _typeface(rPr, "a:ea") == "Yu Gothic UI"


# ── Smoke test: full save round-trip in Japanese mode ────────────────────────


class TestSaveRoundTrip:
    def test_can_save_japanese_presentation_to_disk(self, tmp_path):
        """End-to-end: build a small slide and ensure the file opens cleanly."""
        pu.set_language("ja")
        prs = pu.create_presentation()
        slide = pu.new_blank_slide(prs)
        pu.add_textbox(
            slide,
            "Azure Container Apps \u306e\u5165\u9580",
            Inches(1),
            Inches(1),
            Inches(11),
            Inches(1),
            font_size=28,
            bold=True,
        )
        pu.add_bullet_list(
            slide,
            [
                "\u30b3\u30f3\u30c6\u30ca\u3092 **\u30b5\u30fc\u30d0\u30fc\u30ec\u30b9** \u3067\u52d5\u304b\u3059",
                "\u30b9\u30b1\u30fc\u30ea\u30f3\u30b0\u3092\u81ea\u52d5\u5316",
            ],
            Inches(1),
            Inches(2.5),
            Inches(11),
        )
        out = tmp_path / "ja_smoke.pptx"
        pu.save_presentation(prs, str(out))
        assert out.exists()
        assert out.stat().st_size > 10_000

        # Re-open and confirm at least one EA typeface is present in slide XML.
        from pptx import Presentation as _P

        reopened = _P(str(out))
        found_ea = False
        for s in reopened.slides:
            for sh in s.shapes:
                if not sh.has_text_frame:
                    continue
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        rPr = _rPr(r)
                        if rPr is not None and rPr.find(qn("a:ea")) is not None:
                            found_ea = True
        assert found_ea, "saved presentation should contain at least one <a:ea> element"
