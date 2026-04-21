"""Tests for pptx_qa_checks.py — slide QA validation."""

import importlib.util
from pathlib import Path
from unittest.mock import MagicMock, PropertyMock

import pytest

# Load pptx_qa_checks from the skill directory (hyphenated dir cannot be
# a normal Python import).
_QA_SCRIPT = (
    Path(__file__).resolve().parent.parent
    / "skills"
    / "pptx-generator"
    / "pptx_qa_checks.py"
)
_spec = importlib.util.spec_from_file_location("pptx_qa_checks", _QA_SCRIPT)
_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_mod)

run_all_checks = _mod.run_all_checks
check_slide_count = _mod.check_slide_count
check_placeholder_text = _mod.check_placeholder_text
check_speaker_notes = _mod.check_speaker_notes
check_font_sizes = _mod.check_font_sizes
check_empty_text_frames = _mod.check_empty_text_frames
check_shape_overflow = _mod.check_shape_overflow
check_slide_text_density = _mod.check_slide_text_density
format_report = _mod.format_report
PLACEHOLDER_RE = _mod.PLACEHOLDER_RE


# ── Helpers to build mock presentations ───────────────────────────────────────


def _make_prs(slide_count: int = 3):
    """Create a mock Presentation with N slides, each having basic shapes."""
    prs = MagicMock()
    slides = []
    for i in range(slide_count):
        slide = MagicMock()
        slide.has_notes_slide = True
        notes_frame = MagicMock()
        notes_frame.text = (
            f"Speaker notes for slide {i + 1} with enough detail to pass checks."
        )
        slide.notes_slide.notes_text_frame = notes_frame

        # One content shape with text
        shape = MagicMock()
        shape.has_text_frame = True
        shape.text_frame.text = f"Content for slide {i + 1}"
        shape.text_frame.paragraphs = []
        shape.left = 914400  # ~1 inch
        shape.top = 914400
        shape.width = 914400 * 8
        shape.height = 914400 * 2
        shape.name = f"TextBox {i}"
        shape._element = MagicMock()
        parent_elem = MagicMock()
        parent_elem.tag = "p:spTree"
        shape._element.getparent.return_value = parent_elem
        # No sub-shapes
        shape.shapes = None
        type(shape).shapes = PropertyMock(side_effect=AttributeError)

        slide.shapes = [shape]
        slides.append(slide)
    prs.slides = slides
    return prs


def _make_shape(text="Hello", left=914400, top=914400, width=914400 * 5, height=914400):
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text_frame.text = text
    shape.text_frame.paragraphs = []
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    shape.name = "TextBox"
    shape._element = MagicMock()
    parent = MagicMock()
    parent.tag = "p:spTree"
    shape._element.getparent.return_value = parent
    return shape


# ── check_slide_count ─────────────────────────────────────────────────────────


class TestSlideCount:
    def test_exact_match(self):
        prs = _make_prs(5)
        issues = check_slide_count(prs, expected=5)
        assert issues == []

    def test_mismatch_minor(self):
        prs = _make_prs(4)
        issues = check_slide_count(prs, expected=5)
        assert len(issues) == 1
        assert issues[0]["severity"] == "MAJOR"

    def test_mismatch_large(self):
        prs = _make_prs(1)
        issues = check_slide_count(prs, expected=10)
        assert len(issues) == 1
        assert issues[0]["severity"] == "CRITICAL"

    def test_zero_slides(self):
        prs = _make_prs(0)
        issues = check_slide_count(prs, expected=None)
        assert any(i["severity"] == "CRITICAL" for i in issues)

    def test_no_expected(self):
        prs = _make_prs(5)
        issues = check_slide_count(prs, expected=None)
        assert issues == []


# ── check_placeholder_text ────────────────────────────────────────────────────


class TestPlaceholderText:
    def test_no_placeholders(self):
        prs = _make_prs(1)
        issues = check_placeholder_text(prs)
        assert issues == []

    def test_detects_todo(self):
        prs = _make_prs(1)
        prs.slides[0].shapes[0].text_frame.text = "TODO: finish this"
        issues = check_placeholder_text(prs)
        assert len(issues) >= 1
        assert issues[0]["check"] == "placeholder_text"
        assert issues[0]["severity"] == "CRITICAL"

    def test_detects_lorem(self):
        prs = _make_prs(1)
        prs.slides[0].shapes[0].text_frame.text = "Lorem ipsum dolor sit amet"
        issues = check_placeholder_text(prs)
        assert len(issues) >= 1

    def test_detects_placeholder_in_notes(self):
        prs = _make_prs(1)
        prs.slides[0].notes_slide.notes_text_frame.text = "FIXME add speaker notes"
        issues = check_placeholder_text(prs)
        placeholder_in_notes = [
            i for i in issues if i["check"] == "placeholder_in_notes"
        ]
        assert len(placeholder_in_notes) >= 1


# ── check_speaker_notes ──────────────────────────────────────────────────────


class TestSpeakerNotes:
    def test_notes_present(self):
        prs = _make_prs(1)
        issues = check_speaker_notes(prs)
        assert issues == []

    def test_missing_notes(self):
        prs = _make_prs(1)
        prs.slides[0].has_notes_slide = False
        issues = check_speaker_notes(prs)
        assert len(issues) == 1
        assert issues[0]["check"] == "missing_notes"

    def test_short_notes(self):
        prs = _make_prs(1)
        prs.slides[0].notes_slide.notes_text_frame.text = "Hi"
        issues = check_speaker_notes(prs)
        short = [i for i in issues if i["check"] == "short_notes"]
        assert len(short) == 1


# ── check_font_sizes ─────────────────────────────────────────────────────────


class TestFontSizes:
    def _make_prs_with_font(self, pt_size):
        prs = _make_prs(1)
        run = MagicMock()
        run.font.size = MagicMock()
        run.font.size.pt = pt_size
        run.text = "Sample text"
        para = MagicMock()
        para.runs = [run]
        prs.slides[0].shapes[0].text_frame.paragraphs = [para]
        return prs

    def test_normal_font(self):
        prs = self._make_prs_with_font(14)
        issues = check_font_sizes(prs)
        assert issues == []

    def test_tiny_font(self):
        prs = self._make_prs_with_font(6)
        issues = check_font_sizes(prs)
        assert len(issues) == 1
        assert issues[0]["check"] == "font_too_small"

    def test_huge_font(self):
        prs = self._make_prs_with_font(60)
        issues = check_font_sizes(prs)
        assert len(issues) == 1
        assert issues[0]["check"] == "font_too_large"


# ── check_slide_text_density ─────────────────────────────────────────────────


class TestTextDensity:
    def test_normal_density(self):
        prs = _make_prs(2)
        issues = check_slide_text_density(prs)
        assert issues == []

    def test_empty_slide_flagged(self):
        prs = _make_prs(2)
        prs.slides[1].shapes[0].text_frame.text = ""
        issues = check_slide_text_density(prs)
        empty = [i for i in issues if i["check"] == "empty_slide"]
        assert len(empty) == 1

    def test_text_heavy_slide(self):
        prs = _make_prs(1)
        prs.slides[0].shapes[0].text_frame.text = "x" * 2500
        issues = check_slide_text_density(prs)
        dense = [i for i in issues if i["check"] == "text_density_high"]
        assert len(dense) == 1


# ── run_all_checks ────────────────────────────────────────────────────────────


class TestRunAllChecks:
    def test_nonexistent_file(self, tmp_path):
        report = run_all_checks(str(tmp_path / "nonexistent.pptx"))
        assert report["status"] == "ERROR"

    def test_report_structure(self, tmp_path):
        """run_all_checks on a bad path returns the expected keys."""
        report = run_all_checks(str(tmp_path / "bad.pptx"))
        assert "status" in report
        assert "issues" in report
        assert "summary" in report
        assert isinstance(report["summary"], dict)

    def test_with_real_minimal_pptx(self, tmp_path):
        """Create a real minimal pptx and run checks against it."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Add a text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        tf.text = "Test slide content"

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            "These are the speaker notes for this slide. They contain enough "
            "text to pass the minimum length check in the QA validator."
        )

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        report = run_all_checks(str(pptx_path), expected_slides=1)
        assert report["status"] in ("CLEAN", "ISSUES_FOUND")
        assert report["slide_count"] == 1
        assert "summary" in report


# ── format_report ─────────────────────────────────────────────────────────────


class TestFormatReport:
    def test_clean_report(self):
        report = {
            "status": "CLEAN",
            "file": "test.pptx",
            "slide_count": 3,
            "expected_slides": 3,
            "issues": [],
            "issues_by_slide": {},
            "summary": {"CRITICAL": 0, "MAJOR": 0, "MINOR": 0},
        }
        text = format_report(report)
        assert "CLEAN" in text
        assert "test.pptx" in text

    def test_issues_report(self):
        report = {
            "status": "ISSUES_FOUND",
            "file": "test.pptx",
            "slide_count": 1,
            "expected_slides": 5,
            "issues": [
                {
                    "slide": 0,
                    "severity": "CRITICAL",
                    "check": "slide_count",
                    "message": "wrong count",
                }
            ],
            "issues_by_slide": {
                0: [
                    {
                        "slide": 0,
                        "severity": "CRITICAL",
                        "check": "slide_count",
                        "message": "wrong count",
                    }
                ]
            },
            "summary": {"CRITICAL": 1, "MAJOR": 0, "MINOR": 0},
        }
        text = format_report(report)
        assert "ISSUES_FOUND" in text
        assert "CRITICAL" in text


# ── PLACEHOLDER_RE ────────────────────────────────────────────────────────────


class TestPlaceholderRegex:
    @pytest.mark.parametrize(
        "text",
        [
            "xxxx",
            "lorem ipsum",
            "TODO",
            "FIXME",
            "placeholder",
            "TBD",
            "insert here",
        ],
    )
    def test_matches_known_patterns(self, text):
        assert PLACEHOLDER_RE.search(text)

    def test_no_match_normal_text(self):
        assert PLACEHOLDER_RE.search("Azure Kubernetes Service overview") is None
